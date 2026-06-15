import os
import logging
import re
from PIL import Image
import PyPDF2
from pptx import Presentation

from backend.config.settings import settings
from backend.utils.exceptions import APIError


def clean_extracted_text(text: str) -> str:
    """Cleans up extracted text from documents:
    - Removes duplicate paragraphs (running headers/footers)
    - Fixes broken line breaks (replaces single newlines with spaces)
    - Normalizes spacing and double paragraph boundaries
    """
    if not text:
        return ""
    
    # Standardize line endings
    normalized = text.replace('\r\n', '\n').replace('\r', '\n')
    
    # Deduplicate paragraphs
    lines = normalized.split('\n')
    unique_lines = []
    seen = set()
    for line in lines:
        line_stripped = line.strip()
        if not line_stripped:
            unique_lines.append("")
            continue
        # Remove duplicate footers/headers or copy-paste duplicates
        if line_stripped in seen:
            continue
        seen.add(line_stripped)
        unique_lines.append(line)
        
    cleaned = '\n'.join(unique_lines)
    
    # Join single newlines with spaces (broken wrap-around lines)
    # A single newline not preceded or followed by another newline is replaced with a space
    cleaned = re.sub(r'(?<!\n)\n(?!\n)', ' ', cleaned)
    
    # Normalize multiple whitespaces
    cleaned = re.sub(r' +', ' ', cleaned)
    
    # Keep max double newlines for paragraph structures
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    
    return cleaned.strip()


def validate_file(file) -> str:
    file_ext = os.path.splitext(file.filename)[1].lower()
    if file_ext not in settings.ALLOWED_FILE_EXTENSIONS:
        raise APIError(
            f"Unsupported file type. Allowed: {', '.join(sorted(settings.ALLOWED_FILE_EXTENSIONS))}",
            status_code=400,
        )
    return file_ext


def extract_text_from_file(file):
    file_ext = validate_file(file)

    if file_ext == '.pdf':
        return extract_text_from_pdf(file)
    if file_ext == '.pptx':
        return extract_text_from_pptx(file)
    if file_ext == '.txt':
        text = file.read().decode('utf-8')
        return clean_extracted_text(text)
    if file_ext in {'.png', '.jpg', '.jpeg', '.webp'}:
        return extract_image(file)

    raise APIError('Unsupported file format', status_code=400)


def extract_text_from_pdf(file):
    try:
        reader = PyPDF2.PdfReader(file)
        text = ''.join((page.extract_text() or '') + '\n' for page in reader.pages)
        cleaned = clean_extracted_text(text)
        
        # Log extracted length
        logging.info(f"[FileService] Extracted PDF text length: {len(cleaned)} characters.")
        
        # Verify extracted text is not empty (e.g. scanned PDF)
        if not cleaned:
            logging.warning("[FileService] Extracted PDF text is empty. PDF might be scanned or image-only.")
            raise APIError(
                "The uploaded PDF contains no extractable text. It appears to be a scanned document or image. "
                "Please upload a text-based PDF or convert it using OCR.",
                status_code=400
            )
        return cleaned
    except APIError:
        raise
    except Exception as exc:
        logging.error('PDF extraction failed', exc_info=True)
        raise APIError(f'Could not extract text from PDF: {str(exc)}', status_code=400)


def extract_text_from_pptx(file):
    try:
        presentation = Presentation(file)
        text = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        cleaned = clean_extracted_text(text)
        logging.info(f"[FileService] Extracted PPTX text length: {len(cleaned)} characters.")
        
        if not cleaned:
            raise APIError("The uploaded PowerPoint presentation does not contain any extractable text.", status_code=400)
            
        return cleaned
    except APIError:
        raise
    except Exception as exc:
        logging.error('PPTX extraction failed', exc_info=True)
        raise APIError(f'Could not extract text from PPTX: {str(exc)}', status_code=400)


def extract_image(file):
    try:
        image = Image.open(file)
        return image
    except Exception as exc:
        logging.error('Image extraction failed', exc_info=True)
        raise APIError(f'Failed to process uploaded image: {str(exc)}', status_code=400)

